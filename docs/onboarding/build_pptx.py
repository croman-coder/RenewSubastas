#!/usr/bin/env python3
"""Render the onboarding decks to editable .pptx files.

Single source of truth: content/*.json (shared with build.mjs).
Brand mirrored from the transactional email system: ink on warm-neutral
paper, semantic color for status only.

Usage: python3 docs/onboarding/build_pptx.py
"""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
CONTENT = ROOT / "content"
ASSETS = ROOT / "assets"
OUT = ROOT / "pptx"
OUT.mkdir(exist_ok=True)

# 16:9 at 13.333 x 7.5 in
EMU_IN = 914400
SW = int(13.333 * EMU_IN)
SH = int(7.5 * EMU_IN)

# palette
INK = RGBColor(0x0A, 0x0A, 0x0A)
PAPER = RGBColor(0xE9, 0xE9, 0xEB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
T2 = RGBColor(0x52, 0x52, 0x5B)
MUTED = RGBColor(0x71, 0x71, 0x7A)
FAINT = RGBColor(0xA1, 0xA1, 0xAA)
HAIR = RGBColor(0xE4, 0xE4, 0xE7)
DGRAY = RGBColor(0xD4, 0xD4, 0xD8)

TONES = {
    "success": (RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x16, 0x65, 0x34)),
    "danger": (RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0x99, 0x1B, 0x1B)),
    "warning": (RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E)),
    "info": (RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1E, 0x40, 0xAF)),
    "neutral": (RGBColor(0x27, 0x27, 0x2A), RGBColor(0xFA, 0xFA, 0xFA)),
}

FONT = "Calibri"  # close system-sans analogue available in PowerPoint
MARGIN = Emu(int(0.92 * EMU_IN))
CONTENT_W = SW - 2 * MARGIN

WORDMARK_DARK = str(ASSETS / "renew-wordmark-black.png")
WORDMARK_LIGHT = str(ASSETS / "renew-wordmark-white.png")
WM_RATIO = 385 / 1080  # h/w


def _strip(s):
    return re.sub(r"\*\*(.+?)\*\*", r"\1", str(s or ""))


def _bold_runs(p, s, base_color, base_size, bold_default=False, font=FONT):
    """Add runs to paragraph p, honoring **bold** spans."""
    parts = re.split(r"(\*\*.+?\*\*)", str(s or ""))
    for part in parts:
        if not part:
            continue
        b = part.startswith("**") and part.endswith("**")
        text = part[2:-2] if b else part
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = base_size
        r.font.color.rgb = base_color
        r.font.bold = b or bold_default


def add_bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    # send to back
    sp = s._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return s


def rrect(slide, x, y, w, h, fill, radius=0.14, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def run(p, text, size, color, bold=False, caps=False, spacing=None, font=FONT):
    r = p.add_run()
    r.text = text.upper() if caps else text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    if spacing is not None:
        _set_spacing(r, spacing)
    return r


def _set_spacing(r, pts):
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def eyebrow_pill(slide, x, y, text):
    w = Emu(int((0.16 * len(text) + 0.5) * EMU_IN))
    w = min(w, Emu(int(6 * EMU_IN)))
    h = Emu(int(0.34 * EMU_IN))
    pill = rrect(slide, x, y, w, h, INK, radius=0.5)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(int(0.12 * EMU_IN))
    tf.margin_right = Emu(int(0.12 * EMU_IN))
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run(p, text, Pt(9), WHITE, bold=True, caps=True, spacing=1.2)
    return y + h


def foot(slide, deck_short, pg, dark=False):
    wm = WORDMARK_LIGHT if dark else WORDMARK_DARK
    h = Emu(int(0.2 * EMU_IN))
    w = Emu(int(0.2 / WM_RATIO * EMU_IN))
    y = SH - Emu(int(0.55 * EMU_IN))
    slide.shapes.add_picture(wm, MARGIN, y, width=w, height=h)
    tb, tf = txt(slide, SW - MARGIN - Emu(int(5 * EMU_IN)), y - Emu(int(0.02 * EMU_IN)),
                 Emu(int(5 * EMU_IN)), Emu(int(0.3 * EMU_IN)), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run(p, f"{deck_short} \u00b7 {pg}", Pt(8), MUTED, bold=True, spacing=0.4)


# ---------- slide builders ----------

def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_cover(prs, s, deck_short):
    sl = slide_blank(prs)
    add_bg(sl, INK)
    x = MARGIN
    y = Emu(int(2.3 * EMU_IN))
    wmw = Emu(int(2.0 * EMU_IN))
    wmh = Emu(int(2.0 * WM_RATIO * EMU_IN))
    sl.shapes.add_picture(WORDMARK_LIGHT, x, y, width=wmw, height=wmh)
    y += wmh + Emu(int(0.45 * EMU_IN))
    if s.get("eyebrow"):
        tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(0.3 * EMU_IN)))
        run(tf.paragraphs[0], s["eyebrow"], Pt(11), FAINT, bold=True, caps=True, spacing=1.4)
        y += Emu(int(0.45 * EMU_IN))
    tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(1.6 * EMU_IN)))
    run(tf.paragraphs[0], _strip(s["title"]), Pt(46), WHITE, bold=True, spacing=-0.5)
    y += Emu(int(1.5 * EMU_IN))
    if s.get("subtitle"):
        tb, tf = txt(sl, x, y, Emu(int(8.2 * EMU_IN)), Emu(int(1.4 * EMU_IN)))
        p = tf.paragraphs[0]
        p.line_spacing = 1.4
        _bold_runs(p, s["subtitle"], DGRAY, Pt(16))
    # footer line
    tb, tf = txt(sl, x, SH - Emu(int(0.7 * EMU_IN)), CONTENT_W, Emu(int(0.3 * EMU_IN)))
    p = tf.paragraphs[0]
    run(p, deck_short, Pt(9), MUTED, bold=True, spacing=0.6)
    run(p, "    Santa Rosa Automotores \u00b7 Paraguay", Pt(9), RGBColor(0x52, 0x52, 0x5B))
    return sl


def build_divider(prs, s):
    sl = slide_blank(prs)
    add_bg(sl, INK)
    x = MARGIN
    y = Emu(int(2.7 * EMU_IN))
    if s.get("num"):
        tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(0.7 * EMU_IN)))
        run(tf.paragraphs[0], s["num"], Pt(28), RGBColor(0x52, 0x52, 0x5B), bold=True)
        y += Emu(int(0.75 * EMU_IN))
    tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(1.3 * EMU_IN)))
    run(tf.paragraphs[0], _strip(s["title"]), Pt(40), WHITE, bold=True, spacing=-0.5)
    y += Emu(int(1.25 * EMU_IN))
    if s.get("subtitle"):
        tb, tf = txt(sl, x, y, Emu(int(8.5 * EMU_IN)), Emu(int(1 * EMU_IN)))
        p = tf.paragraphs[0]
        p.line_spacing = 1.4
        _bold_runs(p, s["subtitle"], FAINT, Pt(16))
    return sl


def build_closing(prs, s, deck_short):
    sl = slide_blank(prs)
    add_bg(sl, INK)
    x = MARGIN
    y = Emu(int(2.5 * EMU_IN))
    wmw = Emu(int(1.7 * EMU_IN))
    wmh = Emu(int(1.7 * WM_RATIO * EMU_IN))
    sl.shapes.add_picture(WORDMARK_LIGHT, x, y, width=wmw, height=wmh)
    y += wmh + Emu(int(0.5 * EMU_IN))
    tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(1.1 * EMU_IN)))
    run(tf.paragraphs[0], _strip(s["title"]), Pt(34), WHITE, bold=True, spacing=-0.4)
    y += Emu(int(1.0 * EMU_IN))
    for line in s.get("lines", []):
        tb, tf = txt(sl, x, y, Emu(int(8.5 * EMU_IN)), Emu(int(0.45 * EMU_IN)))
        run(tf.paragraphs[0], _strip(line), Pt(14), FAINT)
        y += Emu(int(0.42 * EMU_IN))
    return sl


def content_header(sl, s):
    """Render eyebrow + title + intro. Return current y (top of body)."""
    x = MARGIN
    y = MARGIN
    if s.get("eyebrow"):
        y = eyebrow_pill(sl, x, y, s["eyebrow"]) + Emu(int(0.22 * EMU_IN))
    if s.get("title"):
        tb, tf = txt(sl, x, y, CONTENT_W, Emu(int(0.95 * EMU_IN)))
        run(tf.paragraphs[0], _strip(s["title"]), Pt(30), INK, bold=True, spacing=-0.4)
        y += Emu(int(0.92 * EMU_IN))
    if s.get("intro"):
        tb, tf = txt(sl, x, y, Emu(int(10.5 * EMU_IN)), Emu(int(0.9 * EMU_IN)))
        p = tf.paragraphs[0]
        p.line_spacing = 1.35
        _bold_runs(p, s["intro"], T2, Pt(14))
        # estimate height
        lines = max(1, int(len(_strip(s["intro"])) / 95) + 1)
        y += Emu(int((0.32 * lines + 0.18) * EMU_IN))
    return x, y


def add_callout(sl, x, y, c, w=None):
    w = w or CONTENT_W
    bg, fg = TONES.get(c.get("tone", "neutral"), TONES["neutral"])
    text = _strip(c.get("text", ""))
    title = c.get("title")
    n = len(text)
    lines = max(1, int(n / 105) + 1) + (1 if title else 0)
    h = Emu(int((0.34 * lines + 0.42) * EMU_IN))
    box = rrect(sl, x, y, w, h, bg, radius=0.12)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(int(0.26 * EMU_IN))
    tf.margin_right = Emu(int(0.26 * EMU_IN))
    tf.margin_top = Emu(int(0.16 * EMU_IN))
    tf.margin_bottom = Emu(int(0.16 * EMU_IN))
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    if title:
        run(p, _strip(title) + "  ", Pt(12.5), fg, bold=True)
    _bold_runs(p, text, fg, Pt(12.5))
    return y + h + Emu(int(0.2 * EMU_IN))


def build_content(prs, s, deck_short, pg):
    sl = slide_blank(prs)
    add_bg(sl, PAPER)
    x, y = content_header(sl, s)
    t = s["type"]

    if t in ("toc", "bullets"):
        items = s.get("items", [])
        avail = SH - y - Emu(int(0.9 * EMU_IN))
        rowh = min(Emu(int(0.62 * EMU_IN)), int(avail / max(1, len(items))))
        for i, it in enumerate(items):
            ry = y + i * rowh
            if t == "toc":
                lead = f"{i+1:02d}"
                text = it["text"] if isinstance(it, dict) else it
            else:
                lead = it.get("lead", "") if isinstance(it, dict) else ""
                text = it["text"] if isinstance(it, dict) else it
            tb, tf = txt(sl, x, ry, CONTENT_W, rowh, MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            if t == "toc":
                run(p, lead + "   ", Pt(15), FAINT, bold=True)
                _bold_runs(p, text, T2, Pt(15))
            else:
                if lead:
                    run(p, _strip(lead) + ":  ", Pt(14.5), INK, bold=True)
                _bold_runs(p, text, T2, Pt(14.5))
            # divider line
            ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry + rowh - Emu(900), CONTENT_W, Emu(900))
            ln.fill.solid(); ln.fill.fore_color.rgb = HAIR; ln.line.fill.background(); ln.shadow.inherit = False

    elif t == "steps":
        steps = s.get("steps", [])
        avail = SH - y - Emu(int(0.9 * EMU_IN))
        rowh = min(Emu(int(1.05 * EMU_IN)), int(avail / max(1, len(steps))))
        d = Emu(int(0.5 * EMU_IN))
        for i, st in enumerate(steps):
            ry = y + i * rowh
            circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x, ry, d, d)
            circ.fill.solid(); circ.fill.fore_color.rgb = INK; circ.line.fill.background(); circ.shadow.inherit = False
            ctf = circ.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            run(cp, str(i + 1), Pt(15), WHITE, bold=True)
            tx = x + d + Emu(int(0.28 * EMU_IN))
            tw = CONTENT_W - d - Emu(int(0.28 * EMU_IN))
            tb, tf = txt(sl, tx, ry - Emu(int(0.02 * EMU_IN)), tw, rowh, MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            run(p, _strip(st["title"]), Pt(15), INK, bold=True)
            p2 = tf.add_paragraph(); p2.line_spacing = 1.28
            _bold_runs(p2, st["text"], T2, Pt(12.5))

    elif t == "statpair":
        tiles = s.get("tiles", [])
        gap = Emu(int(0.25 * EMU_IN))
        tw = int((CONTENT_W - gap * (len(tiles) - 1)) / len(tiles))
        th = Emu(int(1.5 * EMU_IN))
        for i, tl in enumerate(tiles):
            tx = x + i * (tw + gap)
            strong = tl.get("strong")
            box = rrect(sl, tx, y, tw, th, INK if strong else SOFT, radius=0.1)
            tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Emu(int(0.34 * EMU_IN)); tf.margin_right = Emu(int(0.3 * EMU_IN))
            p = tf.paragraphs[0]
            run(p, tl["label"], Pt(10.5), FAINT if strong else MUTED, bold=True, caps=True, spacing=1.0)
            p2 = tf.add_paragraph()
            run(p2, tl["value"], Pt(24), WHITE if strong else INK, bold=True, spacing=-0.3)
        y += th + Emu(int(0.25 * EMU_IN))
        if s.get("callout"):
            add_callout(sl, x, y, s["callout"])

    elif t == "kv":
        rows = s.get("rows", [])
        rowh = Emu(int(0.5 * EMU_IN))
        kw = int(CONTENT_W * 0.36)
        for i, (k, v) in enumerate(rows):
            ry = y + i * rowh
            tb, tf = txt(sl, x, ry, kw, rowh, MSO_ANCHOR.MIDDLE)
            _bold_runs(tf.paragraphs[0], k, MUTED, Pt(13.5))
            tb2, tf2 = txt(sl, x + kw, ry, CONTENT_W - kw, rowh, MSO_ANCHOR.MIDDLE)
            _bold_runs(tf2.paragraphs[0], v or "\u2014", INK, Pt(13.5), bold_default=True)
            ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry + rowh - Emu(900), CONTENT_W, Emu(900))
            ln.fill.solid(); ln.fill.fore_color.rgb = HAIR; ln.line.fill.background(); ln.shadow.inherit = False

    elif t == "table":
        build_table(sl, s, x, y)

    elif t == "flow":
        build_flow(sl, s, x, y)

    foot(sl, deck_short, pg)
    return sl


def build_table(sl, s, x, y):
    head = s.get("head", [])
    rows = s.get("rows", [])
    ncol = len(head)
    # column weights: first col narrower-ish, rest share
    if ncol == 3:
        weights = [0.22, 0.56, 0.22]
    elif ncol == 2:
        weights = [0.32, 0.68]
    else:
        weights = [1.0 / ncol] * ncol
    cols = [int(CONTENT_W * w) for w in weights]
    callout = s.get("callout")
    bottom_reserve = Emu(int(1.5 * EMU_IN)) if callout else Emu(int(0.9 * EMU_IN))
    avail = SH - y - bottom_reserve
    headh = Emu(int(0.5 * EMU_IN))
    rowh = min(Emu(int(0.72 * EMU_IN)), int((avail - headh) / max(1, len(rows))))

    # header
    cx = x
    for j, htext in enumerate(head):
        tb, tf = txt(sl, cx, y, cols[j], headh, MSO_ANCHOR.BOTTOM)
        tf.margin_right = Emu(int(0.12 * EMU_IN))
        run(tf.paragraphs[0], htext, Pt(10.5), MUTED, bold=True, caps=True, spacing=0.6)
        cx += cols[j]
    # header underline (ink, 2px)
    ul = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + headh, CONTENT_W, Emu(1800))
    ul.fill.solid(); ul.fill.fore_color.rgb = INK; ul.line.fill.background(); ul.shadow.inherit = False

    ry0 = y + headh + Emu(int(0.06 * EMU_IN))
    for i, row in enumerate(rows):
        ry = ry0 + i * rowh
        cx = x
        for j, cell in enumerate(row):
            tb, tf = txt(sl, cx, ry, cols[j], rowh, MSO_ANCHOR.MIDDLE)
            tf.margin_right = Emu(int(0.18 * EMU_IN))
            p = tf.paragraphs[0]; p.line_spacing = 1.2
            if j == 0:
                _bold_runs(p, cell, INK, Pt(12.5), bold_default=True)
            else:
                _bold_runs(p, cell, T2, Pt(12))
            cx += cols[j]
        ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry + rowh - Emu(900), CONTENT_W, Emu(900))
        ln.fill.solid(); ln.fill.fore_color.rgb = HAIR; ln.line.fill.background(); ln.shadow.inherit = False

    if callout:
        cy = ry0 + len(rows) * rowh + Emu(int(0.14 * EMU_IN))
        add_callout(sl, x, cy, callout)


def build_flow(sl, s, x, y):
    stages = s.get("stages", [])
    callout = s.get("callout")
    n = len(stages)
    arrow = Emu(int(0.38 * EMU_IN))
    total_arrows = arrow * (n - 1)
    cardw = int((CONTENT_W - total_arrows) / n)
    cardh = Emu(int(1.85 * EMU_IN))
    cy = y + Emu(int(0.1 * EMU_IN))
    cx = x
    for i, st in enumerate(stages):
        box = rrect(sl, cx, cy, cardw, cardh, CARD, radius=0.08, line=HAIR, line_w=Pt(1))
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP; tf.word_wrap = True
        tf.margin_left = Emu(int(0.2 * EMU_IN)); tf.margin_right = Emu(int(0.2 * EMU_IN))
        tf.margin_top = Emu(int(0.2 * EMU_IN))
        p = tf.paragraphs[0]; p.line_spacing = 1.1
        run(p, _strip(st["label"]), Pt(14), INK, bold=True)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.25; p2.space_before = Pt(5)
        _bold_runs(p2, st["sub"], MUTED, Pt(10.5))
        if i < n - 1:
            ax = cx + cardw
            tb, taf = txt(sl, ax, cy, arrow, cardh, MSO_ANCHOR.MIDDLE)
            ap = taf.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
            run(ap, "\u203a", Pt(22), FAINT, bold=True)
        cx += cardw + arrow
    if callout:
        add_callout(sl, x, cy + cardh + Emu(int(0.22 * EMU_IN)), callout)


def render_deck(path):
    deck = json.loads(path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    short = deck.get("short", deck["title"])
    pg = 0
    for s in deck["slides"]:
        t = s["type"]
        if t == "cover":
            build_cover(prs, s, short)
        elif t == "divider":
            build_divider(prs, s)
        elif t == "closing":
            build_closing(prs, s, short)
        else:
            pg += 1
            build_content(prs, s, short, pg)
    out = OUT / (path.stem + ".pptx")
    prs.save(str(out))
    print(f"built pptx/{out.name}  ({len(deck['slides'])} slides)")


def main():
    for f in sorted(CONTENT.glob("*.json")):
        render_deck(f)


if __name__ == "__main__":
    main()
